from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.parsers.base import looks_like_equation
from app.schemas.parsed_document import (
    DocumentMetadata,
    EquationRef,
    FigureRef,
    ParsedDocument,
    Section,
    TableBlock,
)


def parse_pptx(file_path: str) -> ParsedDocument:
    prs = Presentation(file_path)
    sections: list[Section] = []
    tables: list[TableBlock] = []
    figures: list[FigureRef] = []
    equations: list[EquationRef] = []
    slide_count = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_count = slide_index
        heading: Optional[str] = None
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if not text.strip():
                    continue
                if shape == slide.shapes.title:
                    heading = text.strip()
                else:
                    body_lines.append(text)
                    for line in text.splitlines():
                        if looks_like_equation(line):
                            equations.append(EquationRef(page=slide_index, text=line.strip()))
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                tables.append(TableBlock(page=slide_index, rows=rows))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                figures.append(FigureRef(page=slide_index))
        sections.append(Section(heading=heading, page=slide_index, text="\n".join(body_lines).strip()))

    return ParsedDocument(
        metadata=DocumentMetadata(source_filename=Path(file_path).name, format="pptx", page_count=slide_count),
        sections=sections,
        tables=tables,
        figures=figures,
        equations=equations,
    )
