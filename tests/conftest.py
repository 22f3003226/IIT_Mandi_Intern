from pathlib import Path

import pytest
from docx import Document as DocxDocument
from fpdf import FPDF
from pptx import Presentation


@pytest.fixture
def sample_txt(tmp_path: Path) -> Path:
    content = (
        "Introduction\n"
        "This chapter introduces Newtons Laws of Motion.\n\n"
        "Newtons First Law\n"
        "An object at rest stays at rest unless acted upon by a force.\n"
        "F = m * a\n"
    )
    path = tmp_path / "sample.txt"
    path.write_text(content, encoding="utf-8")
    return path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    doc = DocxDocument()
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph("This chapter introduces Newtons Laws of Motion.")
    doc.add_heading("Newtons First Law", level=1)
    doc.add_paragraph("An object at rest stays at rest unless acted upon by a force.")
    doc.add_paragraph("F = m * a")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Definition"
    table.cell(1, 0).text = "Force"
    table.cell(1, 1).text = "A push or pull"
    path = tmp_path / "sample.docx"
    doc.save(path)
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Newtons Laws"
    slide.placeholders[1].text = "An object at rest stays at rest unless acted upon by a force.\nF = m * a"
    path = tmp_path / "sample.pptx"
    prs.save(path)
    return path


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=14)
    pdf.cell(0, 10, "Introduction", ln=True)
    pdf.set_font("Helvetica", size=11)
    pdf.multi_cell(0, 8, "This chapter introduces Newtons Laws of Motion.")
    pdf.cell(0, 10, "Newtons First Law", ln=True)
    pdf.multi_cell(0, 8, "An object at rest stays at rest unless acted upon by a force.")
    pdf.cell(0, 8, "F = m * a", ln=True)
    path = tmp_path / "sample.pdf"
    pdf.output(str(path))
    return path


@pytest.fixture
def blank_scanned_pdf(tmp_path: Path) -> Path:
    pdf = FPDF()
    pdf.add_page()
    path = tmp_path / "scanned.pdf"
    pdf.output(str(path))
    return path
